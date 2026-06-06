import os
import re
import fitz  # PyMuPDF
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from docx import Document
from docx.text.paragraph import Paragraph
from docx.table import Table
import openpyxl
from html.parser import HTMLParser

# Helper function to generate Markdown tables from 2D arrays
def make_markdown_table(rows):
    if not rows:
        return ""
    
    # Clean up empty columns or rows where everything is blank
    col_has_data = [False] * len(rows[0])
    for row in rows:
        for idx, cell in enumerate(row):
            if cell is not None and str(cell).strip():
                col_has_data[idx] = True
                
    # Filter columns that contain at least some data
    valid_col_indices = [idx for idx, has_data in enumerate(col_has_data) if has_data]
    if not valid_col_indices:
        return ""
        
    markdown_rows = []
    
    # Process header
    raw_headers = rows[0]
    headers = [str(raw_headers[idx] or "").replace("\n", " ").replace("|", "\\|").strip() for idx in valid_col_indices]
    # If header is completely empty, use Column 1, Column 2, etc.
    headers = [header if header else f"Col {i+1}" for i, header in enumerate(headers)]
    
    markdown_rows.append("| " + " | ".join(headers) + " |")
    markdown_rows.append("| " + " | ".join(["---"] * len(headers)) + " |")
    
    # Process data rows
    for row in rows[1:]:
        cells = []
        for idx in valid_col_indices:
            val = row[idx] if idx < len(row) else ""
            cells.append(str(val or "").replace("\n", " ").replace("|", "\\|").strip())
        markdown_rows.append("| " + " | ".join(cells) + " |")
        
    return "\n" + "\n".join(markdown_rows) + "\n"

# Helper check for rectangle overlaps
def is_inside(rect_a, rect_b):
    # Checks if rect_a (x0, y0, x1, y1) overlaps significantly with rect_b (x0, y0, x1, y1)
    ax0, ay0, ax1, ay1 = rect_a
    bx0, by0, bx1, by1 = rect_b
    
    overlap_x = max(0, min(ax1, bx1) - max(ax0, bx0))
    overlap_y = max(0, min(ay1, by1) - max(ay0, by0))
    overlap_area = overlap_x * overlap_y
    
    area_a = (ax1 - ax0) * (ay1 - ay0)
    if area_a <= 0:
        return False
    return (overlap_area / area_a) > 0.5

# 1. PDF Parser using PyMuPDF
def convert_pdf(pdf_path, output_dir, extract_images=True, extract_tables=True):
    doc = fitz.open(pdf_path)
    markdown_content = []
    
    images_dir = os.path.join(output_dir, "images")
    if extract_images:
        os.makedirs(images_dir, exist_ok=True)
        
    global_image_counter = 1
    
    for page_num in range(len(doc)):
        page = doc[page_num]
        markdown_content.append(f"\n<!-- Page {page_num + 1} -->\n")
        
        # 1. Gather all tables
        tables = []
        table_bboxes = []
        if extract_tables:
            try:
                found_tabs = page.find_tables()
                for tab in found_tabs:
                    tables.append(tab)
                    table_bboxes.append(tab.bbox)
            except Exception as e:
                print(f"Error finding tables on page {page_num+1}: {e}")
                
        # 2. Gather image boxes
        image_items = []
        try:
            image_info = page.get_image_info(xrefs=True)
            for img in image_info:
                bbox = img.get("bbox")
                xref = img.get("xref")
                if bbox and xref:
                    image_items.append({"bbox": bbox, "xref": xref})
        except Exception as e:
            print(f"Error getting image info on page {page_num+1}: {e}")
            
        # 3. Gather text blocks and filter out those inside tables
        text_items = []
        try:
            blocks = page.get_text("blocks")
            for block in blocks:
                bbox = block[:4]
                text = block[4].strip()
                if not text:
                    continue
                    
                # Skip if text overlaps with a table
                inside_table = False
                for t_bbox in table_bboxes:
                    if is_inside(bbox, t_bbox):
                        inside_table = True
                        break
                if not inside_table:
                    text_items.append({"bbox": bbox, "text": text})
        except Exception as e:
            print(f"Error getting text blocks on page {page_num+1}: {e}")
            
        # 4. Combine elements for grid-based visual flow sorting
        page_elements = []
        
        # Add tables
        for idx, tab in enumerate(tables):
            page_elements.append({
                "type": "table",
                "bbox": tab.bbox,
                "data": tab.extract()
            })
            
        # Add text blocks
        for t_item in text_items:
            page_elements.append({
                "type": "text",
                "bbox": t_item["bbox"],
                "content": t_item["text"]
            })
            
        # Add images
        for idx, img_item in enumerate(image_items):
            # Check if image overlaps with text/tables (avoid raw text overlays)
            page_elements.append({
                "type": "image",
                "bbox": img_item["bbox"],
                "xref": img_item["xref"]
            })
            
        # Sort page elements by top coordinate (with a 15px grid-row height allowance), then left coordinate
        page_elements.sort(key=lambda item: (round(item["bbox"][1] / 15) * 15, item["bbox"][0]))
        
        # Render page elements to Markdown
        for el in page_elements:
            if el["type"] == "text":
                # Clean clean headers/paragraphs
                content = el["content"]
                # Convert PDF bold fonts markers if they look like headers (brief single lines)
                if len(content) < 100 and content.isupper() and not content.endswith(('.', ':', ';')):
                    markdown_content.append(f"\n### {content}\n")
                else:
                    markdown_content.append(f"\n{content}\n")
                    
            elif el["type"] == "table":
                md_table = make_markdown_table(el["data"])
                markdown_content.append(f"\n{md_table}\n")
                
            elif el["type"] == "image" and extract_images:
                xref = el["xref"]
                try:
                    base_image = doc.extract_image(xref)
                    if base_image:
                        image_bytes = base_image["image"]
                        ext = base_image["ext"]
                        img_filename = f"image_{global_image_counter}.{ext}"
                        img_path = os.path.join(images_dir, img_filename)
                        
                        with open(img_path, "wb") as f:
                            f.write(image_bytes)
                            
                        markdown_content.append(f"\n![Image {global_image_counter}](images/{img_filename})\n")
                        global_image_counter += 1
                except Exception as e:
                    print(f"Error rendering image xref {xref} on page {page_num+1}: {e}")
                    
    doc.close()
    return "".join(markdown_content)

# 2. PowerPoint Parser using python-pptx
def convert_pptx(pptx_path, output_dir, extract_images=True):
    try:
        prs = Presentation(pptx_path)
    except Exception as e:
        return f"Error loading PowerPoint file: {str(e)}"
        
    markdown_content = []
    images_dir = os.path.join(output_dir, "images")
    if extract_images:
        os.makedirs(images_dir, exist_ok=True)
        
    global_image_counter = 1
    
    # Helper to unpack grouped shapes
    def unpack_shapes(shapes_list):
        unpacked = []
        for shape in shapes_list:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                try:
                    unpacked.extend(unpack_shapes(shape.shapes))
                except Exception:
                    pass
            else:
                unpacked.append(shape)
        return unpacked
        
    for slide_idx, slide in enumerate(prs.slides):
        markdown_content.append(f"\n## Slide {slide_idx + 1}\n")
        
        shapes = unpack_shapes(slide.shapes)
        slide_elements = []
        
        # Slide height for sorting tolerance (approx 5% of height)
        tolerance = prs.slide_height * 0.05 if hasattr(prs, "slide_height") else 500000
        
        for shape in shapes:
            left, top = shape.left, shape.top
            
            # Text Frame shapes
            if shape.has_text_frame:
                paragraphs_md = []
                for p in shape.text_frame.paragraphs:
                    text = p.text.strip()
                    if not text:
                        continue
                    level = p.level
                    indent = "  " * level
                    if level > 0:
                        paragraphs_md.append(f"{indent}- {text}")
                    else:
                        # Headers or normal lines
                        paragraphs_md.append(text)
                
                content = "\n".join(paragraphs_md)
                if content:
                    slide_elements.append({
                        "type": "text",
                        "top": top,
                        "left": left,
                        "content": content
                    })
                    
            # Tables
            elif shape.has_table:
                table_data = []
                for row in shape.table.rows:
                    row_cells = [cell.text.strip() for cell in row.cells]
                    table_data.append(row_cells)
                slide_elements.append({
                    "type": "table",
                    "top": top,
                    "left": left,
                    "data": table_data
                })
                
            # Pictures
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE and extract_images:
                try:
                    image = shape.image
                    slide_elements.append({
                        "type": "image",
                        "top": top,
                        "left": left,
                        "bytes": image.blob,
                        "ext": image.ext
                    })
                except Exception as e:
                    print(f"Error fetching picture shape from slide {slide_idx+1}: {e}")
                    
        # Sort shapes top-to-bottom, left-to-right
        slide_elements.sort(key=lambda s: (round(s["top"] / tolerance) * tolerance, s["left"]))
        
        for el in slide_elements:
            if el["type"] == "text":
                markdown_content.append(f"\n{el['content']}\n")
                
            elif el["type"] == "table":
                md_table = make_markdown_table(el["data"])
                markdown_content.append(f"\n{md_table}\n")
                
            elif el["type"] == "image":
                try:
                    img_filename = f"slide_{slide_idx+1}_img_{global_image_counter}.{el['ext']}"
                    img_path = os.path.join(images_dir, img_filename)
                    with open(img_path, "wb") as f:
                        f.write(el["bytes"])
                    markdown_content.append(f"\n![Slide Image {global_image_counter}](images/{img_filename})\n")
                    global_image_counter += 1
                except Exception as e:
                    print(f"Error exporting image on slide {slide_idx+1}: {e}")
                    
    return "".join(markdown_content)

# 3. Word Parser using python-docx
def convert_docx(docx_path, output_dir, extract_images=True):
    try:
        doc = Document(docx_path)
    except Exception as e:
        return f"Error loading Word document: {str(e)}"
        
    markdown_content = []
    
    # Extract order preserved body elements
    for child in doc.element.body:
        if child.tag.endswith('p'):
            p = Paragraph(child, doc)
            text = p.text.strip()
            if not text:
                continue
            
            style_name = p.style.name if p.style else ""
            if style_name.startswith('Heading 1'):
                markdown_content.append(f"\n# {text}\n")
            elif style_name.startswith('Heading 2'):
                markdown_content.append(f"\n## {text}\n")
            elif style_name.startswith('Heading 3'):
                markdown_content.append(f"\n### {text}\n")
            elif style_name.startswith('Heading 4'):
                markdown_content.append(f"\n#### {text}\n")
            elif style_name.startswith('List Bullet'):
                markdown_content.append(f"- {text}")
            elif style_name.startswith('List Number'):
                markdown_content.append(f"1. {text}")
            else:
                markdown_content.append(f"\n{text}\n")
                
        elif child.tag.endswith('tbl'):
            t = Table(child, doc)
            rows_data = []
            for row in t.rows:
                row_cells = [cell.text.strip() for cell in row.cells]
                rows_data.append(row_cells)
            md_table = make_markdown_table(rows_data)
            markdown_content.append(f"\n{md_table}\n")
            
    # Extract zipped media images
    if extract_images:
        import zipfile
        images_dir = os.path.join(output_dir, "images")
        os.makedirs(images_dir, exist_ok=True)
        try:
            with zipfile.ZipFile(docx_path) as z:
                img_idx = 1
                for f in z.namelist():
                    if f.startswith('word/media/'):
                        ext = f.split('.')[-1]
                        image_data = z.read(f)
                        img_filename = f"docx_img_{img_idx}.{ext}"
                        with open(os.path.join(images_dir, img_filename), "wb") as img_f:
                            img_f.write(image_data)
                        markdown_content.append(f"\n![Doc Image {img_idx}](images/{img_filename})\n")
                        img_idx += 1
        except Exception as e:
            print(f"Error extracting ZIP media from docx: {e}")
            
    return "\n".join(markdown_content)

# 4. Excel Parser using openpyxl
def convert_xlsx(xlsx_path):
    try:
        wb = openpyxl.load_workbook(xlsx_path, data_only=True)
    except Exception as e:
        return f"Error loading Excel sheets: {str(e)}"
        
    markdown_content = []
    
    for sheet_name in wb.sheetnames:
        sheet = wb[sheet_name]
        markdown_content.append(f"\n## Sheet: {sheet_name}\n")
        
        # Read rows
        rows = list(sheet.iter_rows(values_only=True))
        if not rows:
            markdown_content.append("*Empty Sheet*\n")
            continue
            
        # Filter rows that are completely empty
        cleaned_rows = []
        for r in rows:
            if any(cell is not None for cell in r):
                cleaned_rows.append([str(c) if c is not None else "" for c in r])
                
        if not cleaned_rows:
            markdown_content.append("*Empty Sheet*\n")
            continue
            
        # Pad columns
        max_cols = max(len(r) for r in cleaned_rows)
        for i in range(len(cleaned_rows)):
            if len(cleaned_rows[i]) < max_cols:
                cleaned_rows[i] += [""] * (max_cols - len(cleaned_rows[i]))
                
        md_table = make_markdown_table(cleaned_rows)
        markdown_content.append(f"\n{md_table}\n")
        
    return "".join(markdown_content)

# 5. HTML Parser using built-in HTMLParser
class CustomHTMLToMarkdown(HTMLParser):
    def __init__(self):
        super().__init__()
        self.markdown = []
        self.list_stack = []
        self.in_heading = False
        self.heading_level = 0
        self.in_bold = False
        self.in_italic = False
        self.in_pre = False
        self.href = None
        
    def handle_starttag(self, tag, attrs):
        attrs_dict = dict(attrs)
        if tag in ["h1", "h2", "h3", "h4", "h5", "h6"]:
            self.in_heading = True
            self.heading_level = int(tag[1])
            self.markdown.append("\n\n" + "#" * self.heading_level + " ")
        elif tag == "p":
            self.markdown.append("\n\n")
        elif tag == "br":
            self.markdown.append("\n")
        elif tag in ["strong", "b"]:
            self.in_bold = True
            self.markdown.append("**")
        elif tag in ["em", "i"]:
            self.in_italic = True
            self.markdown.append("*")
        elif tag == "pre":
            self.in_pre = True
            self.markdown.append("\n```\n")
        elif tag == "code":
            if not self.in_pre:
                self.markdown.append("`")
        elif tag == "a":
            self.href = attrs_dict.get("href")
            self.markdown.append("[")
        elif tag in ["ul", "ol"]:
            self.list_stack.append(tag)
            self.markdown.append("\n")
        elif tag == "li":
            indent = "  " * (len(self.list_stack) - 1)
            if self.list_stack and self.list_stack[-1] == "ol":
                self.markdown.append(f"\n{indent}1. ")
            else:
                self.markdown.append(f"\n{indent}- ")
                
    def handle_endtag(self, tag):
        if tag in ["h1", "h2", "h3", "h4", "h5", "h6"]:
            self.in_heading = False
            self.markdown.append("\n\n")
        elif tag == "p":
            self.markdown.append("\n\n")
        elif tag in ["strong", "b"]:
            self.in_bold = False
            self.markdown.append("**")
        elif tag in ["em", "i"]:
            self.in_italic = False
            self.markdown.append("*")
        elif tag == "pre":
            self.in_pre = False
            self.markdown.append("\n```\n")
        elif tag == "code":
            if not self.in_pre:
                self.markdown.append("`")
        elif tag == "a":
            if self.href:
                self.markdown.append(f"]({self.href})")
            else:
                self.markdown.append("]")
            self.href = None
        elif tag in ["ul", "ol"]:
            if self.list_stack:
                self.list_stack.pop()
            self.markdown.append("\n")
            
    def handle_data(self, data):
        if self.in_pre:
            self.markdown.append(data)
        else:
            cleaned_data = re.sub(r'\s+', ' ', data)
            if cleaned_data:
                self.markdown.append(cleaned_data)
                
    def get_markdown(self):
        text = "".join(self.markdown)
        text = re.sub(r'\n{3,}', '\n\n', text)
        return text.strip()

def convert_html(html_path):
    try:
        with open(html_path, "r", encoding="utf-8", errors="ignore") as f:
            html_content = f.read()
    except Exception as e:
        return f"Error reading HTML file: {str(e)}"
        
    parser = CustomHTMLToMarkdown()
    parser.feed(html_content)
    return parser.get_markdown()

# General Dispatcher
def convert_to_markdown(file_path, output_dir, extract_images=True, extract_tables=True):
    filename = os.path.basename(file_path)
    _, ext = os.path.splitext(filename.lower())
    
    if ext == ".pdf":
        return convert_pdf(file_path, output_dir, extract_images, extract_tables)
    elif ext == ".pptx":
        return convert_pptx(file_path, output_dir, extract_images)
    elif ext == ".docx":
        return convert_docx(file_path, output_dir, extract_images)
    elif ext == ".xlsx":
        return convert_xlsx(file_path)
    elif ext in [".html", ".htm"]:
        return convert_html(file_path)
    elif ext == ".txt":
        try:
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                return f.read()
        except Exception as e:
            return f"Error reading text file: {str(e)}"
    else:
        return f"Unsupported file extension: {ext}"
